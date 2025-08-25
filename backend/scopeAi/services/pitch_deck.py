# services/pitch_deck.py
import os
import json
import re
from pathlib import Path
from typing import List
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.util import Pt
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.prompts import ChatPromptTemplate
from langchain.output_parsers import PydanticOutputParser
from .config import MODEL_NAME, SECRET_KEY
import uuid

DEFAULT_MODEL = "gemini-1.5-flash-8b"

# ---- Structured outline from LLM ----
class Slide(BaseModel):
    title: str
    bullets: List[str] = Field(default_factory=list)

class DeckOutline(BaseModel):
    deck_title: str
    subtitle: str | None = ""
    slides: List[Slide]

OUTLINE_TMPL = """You are a pitch deck writer.
Based ONLY on this idea: "{idea}", produce a concise startup pitch deck outline.

Rules:
- 10–12 slides max.
- Each slide: 3–5 short bullets (<= 12 words).
- Cover: Problem, Solution, Product, Market, Business Model,
  Competition, Go-To-Market, Traction (or Early Proof), Financials (or Unit Economics),
  Team, Roadmap, Ask/Use of Funds (adjust to idea).
- No markdown, no code fences, no extra prose.
- Return ONLY valid JSON with this exact schema:

{{
  "deck_title": "string",
  "subtitle": "string", 
  "slides": [
    {{"title": "Problem", "bullets": ["bullet 1", "bullet 2"]}},
    {{"title": "Solution", "bullets": ["bullet 1", "bullet 2"]}}
  ]
}}

IMPORTANT: Return only the JSON object, no extra text before or after."""

def clean_json_response(raw_response: str) -> str:
    """Clean and extract JSON from LLM response."""
    # Remove common prefixes/suffixes that LLMs add
    raw_response = raw_response.strip()
    
    # Remove markdown code blocks if present
    raw_response = re.sub(r'^```(?:json)?\s*', '', raw_response, flags=re.MULTILINE)
    raw_response = re.sub(r'\s*```$', '', raw_response, flags=re.MULTILINE)
    
    # Find JSON object boundaries
    json_match = re.search(r'(\{.*\})', raw_response, re.DOTALL)
    if json_match:
        return json_match.group(1).strip()
    
    return raw_response.strip()

def generate_outline(idea: str, model_name: str = DEFAULT_MODEL) -> DeckOutline:
    if not SECRET_KEY:
        raise RuntimeError("GOOGLE_API_KEY is not set")

    llm = ChatGoogleGenerativeAI(
        model=model_name, 
        temperature=0, 
        google_api_key=SECRET_KEY
    )

    prompt = ChatPromptTemplate.from_template(OUTLINE_TMPL)
    msg = prompt.format_messages(idea=idea)
    raw = llm.invoke(msg).content.strip()
    
    print(f"Raw LLM response: {repr(raw)}")  # Debug logging
    
    try:
        # First, try direct parsing with Pydantic parser
        parser = PydanticOutputParser(pydantic_object=DeckOutline)
        return parser.parse(raw)
    except Exception as e1:
        print(f"Pydantic parser error: {e1}")
        
        try:
            # Clean the response and try manual JSON parsing
            cleaned = clean_json_response(raw)
            print(f"Cleaned response: {repr(cleaned)}")  # Debug logging
            
            parsed_json = json.loads(cleaned)
            return DeckOutline(**parsed_json)
        except Exception as e2:
            print(f"Manual JSON parse error: {e2}")
            
            # Final fallback: create a basic outline
            print("Creating fallback outline...")
            return create_fallback_outline(idea)

def create_fallback_outline(idea: str) -> DeckOutline:
    """Create a basic outline when LLM parsing fails."""
    return DeckOutline(
        deck_title=f"Pitch: {idea[:50]}..." if len(idea) > 50 else f"Pitch: {idea}",
        subtitle="Startup Pitch Deck",
        slides=[
            Slide(title="Problem", bullets=["Market problem exists", "Current solutions insufficient", "Large market opportunity"]),
            Slide(title="Solution", bullets=["Our innovative approach", "Key differentiators", "Proven concept"]),
            Slide(title="Product", bullets=["Core features", "User benefits", "Technical advantages"]),
            Slide(title="Market", bullets=["Target market size", "Growing demand", "Market trends"]),
            Slide(title="Business Model", bullets=["Revenue streams", "Pricing strategy", "Scalable approach"]),
            Slide(title="Competition", bullets=["Competitive landscape", "Our advantages", "Market positioning"]),
            Slide(title="Go-To-Market", bullets=["Customer acquisition", "Sales strategy", "Marketing channels"]),
            Slide(title="Traction", bullets=["Early results", "Customer feedback", "Key metrics"]),
            Slide(title="Financials", bullets=["Revenue projections", "Unit economics", "Path to profitability"]),
            Slide(title="Team", bullets=["Founding team", "Key expertise", "Advisory board"]),
            Slide(title="Ask", bullets=["Funding amount", "Use of funds", "Next milestones"])
        ]
    )

# ---- PPT builder ----
def _ensure_dir(p: Path) -> None:
    p.mkdir(parents=True, exist_ok=True)

def _add_title_slide(prs: Presentation, title: str, subtitle: str | None = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title
    slide.shapes.title.text = (title or "Pitch Deck")[:120]
    if subtitle is not None and len(slide.placeholders) > 1:
        slide.placeholders[1].text = (subtitle or "")[:300]

def _add_bullet_slide(prs: Presentation, heading: str, bullets: List[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title+Content
    slide.shapes.title.text = heading[:120] if heading else "Slide"
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    if not bullets:
        bullets = ["Content not available"]
    for i, b in enumerate(bullets[:6]):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.text = b[:200]  # Truncate very long bullets
        p.level = 0
        p.font.size = Pt(20)

def build_pitch_pptx(deck_id: str, outline: DeckOutline, media_root: Path) -> str:
    """Save to media/pitch/<deck_id>.pptx, return relative path."""
    out_dir = media_root / "pitch"
    _ensure_dir(out_dir)
   # ensure uniqueness
    unique_name = f"{deck_id}_{uuid.uuid4().hex[:8]}.pptx"
    path = out_dir / unique_name

    prs = Presentation()
    _add_title_slide(prs, outline.deck_title, outline.subtitle or "")
    for sl in outline.slides[:12]:
        _add_bullet_slide(prs, sl.title, sl.bullets)
    prs.save(path)
    return path.relative_to(media_root).as_posix()